"""Common primitives for Digital Marketing Catch-Up deck builders."""

from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------- Design tokens ----------
NAVY = RGBColor(0x0F, 0x2A, 0x4F)
NAVY_DARK = RGBColor(0x0A, 0x16, 0x2A)
NAVY_SOFT = RGBColor(0x2B, 0x44, 0x6B)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
ORANGE_DARK = RGBColor(0xB8, 0x5E, 0x1A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
LIGHT_GRAY = RGBColor(0xD7, 0xDD, 0xE8)
MID_GRAY = RGBColor(0x5A, 0x6A, 0x80)
SOFT_GRAY = RGBColor(0x8A, 0x96, 0xA8)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BG = RGBColor(0xFF, 0xF3, 0xE6)
SUCCESS = RGBColor(0x2E, 0x8B, 0x57)
WARN = RGBColor(0xC6, 0x45, 0x2D)
SEA_BG = RGBColor(0xEA, 0xF7, 0xFB)
SEA_PALE = RGBColor(0xD8, 0xF0, 0xF7)
TEAL = RGBColor(0x2E, 0x8F, 0xA3)
TEAL_DARK = RGBColor(0x1E, 0x65, 0x75)

JP_FONT = "Hiragino Sans"
EN_FONT = "Helvetica Neue"


def new_presentation():
    """Create a 16:9 widescreen presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ---------- XML typeface helper ----------
def _set_typefaces(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


# ---------- Primitives ----------
def add_rect(slide, left, top, width, height, *, fill=NAVY, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(shape_type, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(0.75)
    if radius:
        try:
            r.adjustments[0] = 0.06
        except Exception:
            pass
    return r


def _set_transparency(shape, amount):
    try:
        shape.fill.transparency = amount
    except Exception:
        pass


def add_picture_cover(slide, path, l, t, w, h, radius=False):
    """Add an image so it covers the frame without distortion or letterboxing."""
    path = Path(path)
    with Image.open(path) as img:
        img_w, img_h = img.size

    img_ratio = img_w / img_h
    frame_ratio = int(w) / int(h)
    crop_left = crop_right = crop_top = crop_bottom = 0

    if img_ratio > frame_ratio:
        visible_w = img_h * frame_ratio
        crop = max(0, (img_w - visible_w) / img_w / 2)
        crop_left = crop_right = crop
    elif img_ratio < frame_ratio:
        visible_h = img_w / frame_ratio
        crop = max(0, (img_h - visible_h) / img_h / 2)
        crop_top = crop_bottom = crop

    pic = slide.shapes.add_picture(str(path), l, t, w, h)
    pic.crop_left = crop_left
    pic.crop_right = crop_right
    pic.crop_top = crop_top
    pic.crop_bottom = crop_bottom

    if radius:
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        frame.fill.background()
        frame.line.color.rgb = WHITE
        frame.line.width = Pt(1.0)
        try:
            frame.adjustments[0] = 0.04
        except Exception:
            pass
    return pic


def photo_card(slide, path, l, t, w, h, caption=None):
    """Rounded photo card with optional small caption at the bottom."""
    add_rect(slide, l, t, w, h, fill=WHITE, line=LIGHT_GRAY, radius=True)
    pad = Inches(0.08)
    caption_h = Inches(0.36) if caption else 0
    add_picture_cover(slide, path, l + pad, t + pad, w - pad * 2, h - pad * 2, radius=True)
    if caption:
        cap = add_rect(
            slide,
            l + pad,
            t + h - pad - caption_h,
            w - pad * 2,
            caption_h,
            fill=NAVY_DARK,
            radius=True,
        )
        _set_transparency(cap, 18)
        add_text(
            slide,
            caption,
            l + Inches(0.22),
            t + h - Inches(0.38),
            w - Inches(0.44),
            Inches(0.24),
            size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    return None


def hero(slide, path, title, subtitle):
    """Full-bleed image hero with a soft lower text band."""
    prs_w = SLIDE_W
    prs_h = SLIDE_H
    add_picture_cover(slide, path, 0, 0, prs_w, prs_h)

    band_h = Inches(2.15)
    y0 = prs_h - band_h
    for i, alpha in enumerate((42, 32, 22, 12)):
        band = add_rect(
            slide,
            0,
            y0 + Inches(0.18) * i,
            prs_w,
            band_h - Inches(0.18) * i,
            fill=NAVY_DARK,
        )
        _set_transparency(band, alpha)

    add_text(
        slide,
        title,
        Inches(0.75),
        prs_h - Inches(1.72),
        Inches(11.9),
        Inches(0.78),
        size=44,
        bold=True,
        color=WHITE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        subtitle,
        Inches(0.78),
        prs_h - Inches(0.86),
        Inches(11.4),
        Inches(0.42),
        size=18,
        color=SEA_PALE,
        line_spacing=1.0,
    )


def add_shape(slide, shape_type, left, top, width, height, *, fill=NAVY, line=None):
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    return s


def add_text(slide, text, left, top, width, height, *,
             size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=JP_FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        _set_typefaces(r, font)
    return tb


def add_runs(slide, runs, left, top, width, height, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs = list of (text, {size, bold, color, font})."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("size", 12))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", TEXT)
        _set_typefaces(r, opts.get("font", JP_FONT))
    return tb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def page_frame(slide, prs, title, subtitle="", pagenum=""):
    """Top thin navy bar + accent dot + title + subtitle + horizontal rule."""
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), fill=NAVY)
    add_rect(slide, Inches(0.5), Inches(0.38), Inches(0.18), Inches(0.18),
             fill=ORANGE, radius=True)
    add_text(slide, title, Inches(0.78), Inches(0.32), Inches(11.5), Inches(0.55),
             size=22, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, subtitle, Inches(0.78), Inches(0.82), Inches(11.5), Inches(0.35),
                 size=11, color=MID_GRAY)
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.33), Inches(0.015),
             fill=LIGHT_GRAY)
    if pagenum:
        add_text(slide, pagenum, Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3),
                 size=9, color=SOFT_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, prs, text="Digital Marketing Catch-Up 2023→2026", page=""):
    add_text(slide, text, Inches(0.5), Inches(7.1), Inches(10), Inches(0.3),
             size=9, color=SOFT_GRAY)
    if page:
        add_text(slide, page, Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3),
                 size=9, color=SOFT_GRAY, align=PP_ALIGN.RIGHT)


def part_divider(prs, part_num, title, subtitle, summary, chapters):
    """Section divider slide for Part 2-8 opening."""
    s = blank(prs)
    # Full navy BG
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=NAVY)
    # Orange accent vertical bar
    add_rect(s, Inches(0.8), Inches(1.2), Inches(0.15), Inches(5.1), fill=ORANGE)
    # Part label
    add_text(s, f"PART {part_num}", Inches(1.2), Inches(1.2), Inches(4), Inches(0.5),
             size=18, bold=True, color=ORANGE, font=EN_FONT)
    # Big title
    add_text(s, title, Inches(1.2), Inches(1.9), Inches(11), Inches(1.5),
             size=54, bold=True, color=WHITE, line_spacing=1.1)
    # Subtitle
    add_text(s, subtitle, Inches(1.2), Inches(3.6), Inches(11), Inches(0.6),
             size=18, color=LIGHT_GRAY, line_spacing=1.3)
    # Summary rule
    add_rect(s, Inches(1.2), Inches(4.5), Inches(10), Inches(0.015),
             fill=ORANGE)
    # Summary text
    add_text(s, summary, Inches(1.2), Inches(4.7), Inches(11), Inches(0.9),
             size=13, color=LIGHT_GRAY, line_spacing=1.5)
    # Chapters list at bottom-right
    y = Inches(5.9)
    for i, ch in enumerate(chapters):
        add_text(s, f"▸ {ch}", Inches(1.2) + Inches(3.5) * (i % 3),
                 y + Inches(0.4) * (i // 3),
                 Inches(3.3), Inches(0.35),
                 size=11, color=WHITE)
    return s
