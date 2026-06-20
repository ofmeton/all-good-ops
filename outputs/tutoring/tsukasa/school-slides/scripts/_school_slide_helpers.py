"""Shared slide builders for Tsukasa's high school deck."""

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from _deck_lib import (
    LIGHT_GRAY,
    MID_GRAY,
    NAVY,
    SEA_BG,
    SEA_PALE,
    SOFT_GRAY,
    TEAL,
    TEAL_DARK,
    TEXT,
    WHITE,
    add_picture_cover,
    add_rect,
    add_shape,
    add_text,
    blank,
    hero,
    photo_card,
)


def bullet_item(slide, text, top, width=6.35):
    x = Inches(0.72)
    add_shape(slide, MSO_SHAPE.OVAL, x, top + Inches(0.13), Inches(0.20), Inches(0.20), fill=TEAL, line=TEAL)
    add_text(
        slide,
        text,
        Inches(1.28),
        top,
        Inches(width),
        Inches(0.68),
        size=17,
        bold=True,
        color=TEXT,
        line_spacing=1.15,
    )


def add_feature_slide(prs, image_root, school, heading, bullets, commute, source_note, caption="校舎"):
    s = blank(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=SEA_BG)
    add_rect(s, 0, 0, prs.slide_width, Inches(0.10), fill=TEAL)
    add_text(s, heading, Inches(0.72), Inches(0.46), Inches(6.4), Inches(0.42), size=24, bold=True, color=NAVY)
    add_text(s, "写真と一緒に見る、学校の空気", Inches(0.74), Inches(0.96), Inches(5.6), Inches(0.28), size=11, color=MID_GRAY)

    y = Inches(1.62)
    for text in bullets:
        bullet_item(s, text, y)
        y += Inches(1.03)

    add_rect(s, Inches(0.72), Inches(5.15), Inches(6.8), Inches(0.62), fill=WHITE, line=LIGHT_GRAY, radius=True)
    add_text(
        s,
        commute,
        Inches(1.02),
        Inches(5.34),
        Inches(6.2),
        Inches(0.24),
        size=15,
        bold=True,
        color=TEAL_DARK,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    photo_card(
        s,
        image_root / school / "processed" / "wiki_campus_hero.jpg",
        Inches(8.05),
        Inches(1.46),
        Inches(4.62),
        Inches(2.82),
        caption=caption,
    )
    add_text(s, source_note, Inches(0.72), Inches(6.82), Inches(10.8), Inches(0.22), size=8, color=SOFT_GRAY)
    return s


def rich_hook(slide, text, top):
    add_rect(slide, Inches(0.72), top, Inches(11.88), Inches(0.62), fill=WHITE, line=LIGHT_GRAY, radius=True)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(0.98), top + Inches(0.22), Inches(0.16), Inches(0.16), fill=TEAL, line=TEAL)
    add_text(
        slide,
        text,
        Inches(1.32),
        top + Inches(0.08),
        Inches(10.9),
        Inches(0.46),
        size=10.4,
        bold=True,
        color=TEXT,
        line_spacing=0.98,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_rich_feature_slide(prs, heading, image_cards, hooks, commute, source_note):
    s = blank(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=SEA_BG)
    add_rect(s, 0, 0, prs.slide_width, Inches(0.10), fill=TEAL)
    add_text(s, heading, Inches(0.72), Inches(0.46), Inches(6.4), Inches(0.42), size=24, bold=True, color=NAVY)
    add_text(s, "写真と一緒に見る、学校の空気", Inches(0.74), Inches(0.96), Inches(5.6), Inches(0.28), size=11, color=MID_GRAY)

    card_top = Inches(1.28)
    card_h = Inches(1.86)
    if len(image_cards) == 1:
        card_w = Inches(8.00)
        card_h = Inches(2.02)
        positions = [(Inches(0.72), card_top, card_w, card_h)]
        hook_y = Inches(3.48)
    elif len(image_cards) == 2:
        gap = Inches(0.28)
        card_w = Inches(5.80)
        positions = [
            (Inches(0.72), card_top, card_w, card_h),
            (Inches(0.72) + card_w + gap, card_top, card_w, card_h),
        ]
        hook_y = Inches(3.36)
    else:
        gap = Inches(0.24)
        card_w = Inches(3.78)
        positions = [
            (Inches(0.72) + (card_w + gap) * i, card_top, card_w, card_h)
            for i in range(len(image_cards))
        ]
        hook_y = Inches(3.36)

    for (path, caption), (left, top, width, height) in zip(image_cards, positions):
        photo_card(s, path, left, top, width, height, caption=caption)

    y = hook_y
    for text in hooks:
        rich_hook(s, text, y)
        y += Inches(0.64)

    add_rect(s, Inches(0.72), Inches(6.16), Inches(6.8), Inches(0.46), fill=WHITE, line=LIGHT_GRAY, radius=True)
    add_text(
        s,
        commute,
        Inches(1.02),
        Inches(6.30),
        Inches(6.2),
        Inches(0.20),
        size=13,
        bold=True,
        color=TEAL_DARK,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(s, source_note, Inches(0.72), Inches(6.92), Inches(10.8), Inches(0.22), size=8, color=SOFT_GRAY)
    return s


def add_full_bleed_hero(prs, image_path, title, subtitle):
    s = blank(prs)
    hero(s, image_path, title, subtitle)
    return s


def add_contained_hero(prs, image_path, title, subtitle):
    s = blank(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=SEA_BG)
    add_picture_cover(s, image_path, Inches(0.74), Inches(0.45), Inches(11.85), Inches(4.55), radius=True)
    add_rect(s, 0, Inches(4.86), prs.slide_width, Inches(2.64), fill=NAVY)
    add_rect(s, 0, Inches(4.86), prs.slide_width, Inches(0.08), fill=TEAL)
    add_text(s, title, Inches(0.88), Inches(5.34), Inches(11.6), Inches(0.72), size=42, bold=True, color=WHITE, line_spacing=1.0)
    add_text(s, subtitle, Inches(0.9), Inches(6.17), Inches(11.1), Inches(0.38), size=18, color=SEA_PALE, line_spacing=1.0)
    return s
