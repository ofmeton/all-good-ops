#!/usr/bin/env python3
"""Merge all part deck pptx files into a single deck.pptx (245p)."""
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
PARTS = [
    "deck_part1.pptx",
    "deck_part2.pptx",
    "deck_part3a.pptx",
    "deck_part3b.pptx",
    "deck_part3c.pptx",
    "deck_part3d.pptx",
    "deck_part3e.pptx",
    "deck_part4.pptx",
    "deck_part5a.pptx",
    "deck_part5b.pptx",
    "deck_part6.pptx",
    "deck_part7.pptx",
    "deck_part8.pptx",
]
OUTPUT = ROOT / "deck.pptx"


def merge_pptx_files(part_paths, output_path):
    """Merge multiple pptx files by copying slides at the XML level."""
    # Use first file as base
    merged = Presentation(str(part_paths[0]))
    base_count = len(merged.slides)
    print(f"Base: {part_paths[0].name} ({base_count} slides)")

    # For remaining files, copy slides into merged
    for part_path in part_paths[1:]:
        src = Presentation(str(part_path))
        for slide in src.slides:
            # Copy the slide layout from source
            slide_layout = merged.slide_layouts[6]  # blank layout
            new_slide = merged.slides.add_slide(slide_layout)
            # Copy all shapes from source slide to new slide
            for shape in slide.shapes:
                el = shape.element
                new_el = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        print(f"Merged: {part_path.name} ({len(src.slides)} slides) → total {len(merged.slides)}")

    merged.save(str(output_path))
    print(f"\nSaved: {output_path}")
    print(f"Total slides: {len(merged.slides)}")


def main():
    part_paths = [ROOT / p for p in PARTS]
    # Verify all exist
    missing = [p for p in part_paths if not p.exists()]
    if missing:
        print(f"Missing files: {missing}")
        return
    merge_pptx_files(part_paths, OUTPUT)


if __name__ == "__main__":
    main()
