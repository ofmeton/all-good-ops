#!/usr/bin/env python3
"""Build Part 1 (Executive Summary, 10p) of Marketing Catch-Up 2023-2026 deck.

Output: outputs/documents/marketing-catchup-2023-2026/deck_part1.pptx
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "deck_part1.pptx"

# ---------- Design tokens ----------
NAVY = RGBColor(0x0F, 0x2A, 0x4F)
NAVY_DARK = RGBColor(0x0A, 0x16, 0x2A)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
ORANGE_DARK = RGBColor(0xB8, 0x5E, 0x1A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
LIGHT_GRAY = RGBColor(0xD7, 0xDD, 0xE8)
MID_GRAY = RGBColor(0x5A, 0x6A, 0x80)
SOFT_GRAY = RGBColor(0x8A, 0x96, 0xA8)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BG = RGBColor(0xFF, 0xF3, 0xE6)

JP_FONT = "Hiragino Sans"
JP_FONT_BOLD = "Hiragino Sans"

# ---------- XML helper ----------
def _set_typefaces(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


# ---------- Primitives ----------
def add_rect(slide, left, top, width, height, *, fill=NAVY, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(shape_type, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(0.75)
    if radius:
        # Reduce corner radius for modern look
        try:
            r.adjustments[0] = 0.06
        except Exception:
            pass
    return r


def add_text(slide, text, left, top, width, height, *,
             size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=JP_FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        _set_typefaces(r, font)
    return tb


def add_runs(slide, runs, left, top, width, height, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs = list of (text, {size, bold, color, font})."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("size", 12))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", TEXT)
        _set_typefaces(r, opts.get("font", JP_FONT))
    return tb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def page_frame(slide, prs, title, subtitle="", pagenum=""):
    """Common header frame: thin navy bar + title + subtitle + page num."""
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), fill=NAVY)
    # Orange accent dot
    add_rect(slide, Inches(0.5), Inches(0.38), Inches(0.18), Inches(0.18), fill=ORANGE, radius=True)
    add_text(slide, title, Inches(0.78), Inches(0.32), Inches(11.5), Inches(0.55),
             size=22, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, subtitle, Inches(0.78), Inches(0.82), Inches(11.5), Inches(0.35),
                 size=11, color=MID_GRAY)
    # Horizontal rule
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.33), Inches(0.015), fill=LIGHT_GRAY)
    if pagenum:
        add_text(slide, pagenum, Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3),
                 size=9, color=SOFT_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, prs, text="Digital Marketing Catch-Up 2023→2026", page=""):
    add_text(slide, text, Inches(0.5), Inches(7.1), Inches(10), Inches(0.3),
             size=9, color=SOFT_GRAY)
    if page:
        add_text(slide, page, Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3),
                 size=9, color=SOFT_GRAY, align=PP_ALIGN.RIGHT)


# ---------- Slide builders ----------
def slide_01_cover(prs):
    s = blank(prs)
    # full navy background left strip
    add_rect(s, 0, 0, Inches(4.5), prs.slide_height, fill=NAVY)
    # accent orange bar
    add_rect(s, Inches(4.5), 0, Inches(0.12), prs.slide_height, fill=ORANGE)
    # large year arc on navy area
    add_text(s, "2023", Inches(0.6), Inches(0.8), Inches(3.5), Inches(1.0),
             size=48, bold=True, color=WHITE)
    add_text(s, "→", Inches(0.6), Inches(1.7), Inches(3.5), Inches(1.0),
             size=48, bold=True, color=ORANGE)
    add_text(s, "2026", Inches(0.6), Inches(2.6), Inches(3.5), Inches(1.0),
             size=48, bold=True, color=WHITE)
    add_text(s, "2.5 years of blank, closed.", Inches(0.6), Inches(3.9), Inches(3.8), Inches(0.5),
             size=14, color=LIGHT_GRAY, font="Helvetica Neue")
    # Right side content
    add_text(s, "Digital Marketing", Inches(5.0), Inches(1.4), Inches(8.0), Inches(0.9),
             size=40, bold=True, color=NAVY)
    add_text(s, "Catch-Up", Inches(5.0), Inches(2.1), Inches(8.0), Inches(0.9),
             size=40, bold=True, color=NAVY)
    # Tag line
    add_rect(s, Inches(5.0), Inches(3.15), Inches(4.5), Inches(0.05), fill=ORANGE)
    add_text(s, "2023年10月 → 2026年4月\n広告運用実務・9媒体・AI活用の現在地",
             Inches(5.0), Inches(3.35), Inches(8.0), Inches(1.3),
             size=16, color=MID_GRAY)
    # Bottom meta
    add_rect(s, Inches(5.0), Inches(6.2), Inches(7.8), Inches(0.015), fill=LIGHT_GRAY)
    add_text(s, "作成 : 工藤陸 ( 内部資料 )",
             Inches(5.0), Inches(6.35), Inches(8.0), Inches(0.35),
             size=11, color=MID_GRAY)
    add_text(s, "2026-04-25 | Part 1 / 8 | 総 250p",
             Inches(5.0), Inches(6.65), Inches(8.0), Inches(0.35),
             size=11, color=SOFT_GRAY)


def slide_02_how_to_use(prs):
    s = blank(prs)
    page_frame(s, prs, "本書の使い方", "How to read this deck", "P02 / 250")

    # 4 info blocks
    items = [
        ("対象読者", "2023年9月まで広告運用の現場にいた自分\n（Google/Yahoo/Meta/YouTube 実務経験あり）"),
        ("期間", "2023年10月 〜 2026年4月（約2.5年）"),
        ("構成", "8 Part / 約250ページ\nPart 1=要約 / Part 2=市況 / Part 3=媒体別 ほか"),
        ("読み方", "通読: 3〜4時間\n辞書的: Part 3の媒体別を引く\n副業復帰特化: P1 → Part 6 → Part 7"),
    ]
    left = Inches(0.6)
    top = Inches(1.5)
    w = Inches(5.85)
    h = Inches(1.15)
    for i, (k, v) in enumerate(items):
        r = i // 2
        c = i % 2
        x = left + (w + Inches(0.3)) * c
        y = top + (h + Inches(0.15)) * r
        add_rect(s, x, y, w, h, fill=LIGHT, radius=True)
        add_rect(s, x, y, Inches(0.08), h, fill=ORANGE)
        add_text(s, k, x + Inches(0.35), y + Inches(0.15), Inches(2.0), Inches(0.3),
                 size=11, bold=True, color=NAVY)
        add_text(s, v, x + Inches(0.35), y + Inches(0.45), w - Inches(0.45), Inches(0.9),
                 size=11, color=TEXT, line_spacing=1.3)

    # Legend
    add_rect(s, Inches(0.6), Inches(4.2), Inches(12.15), Inches(0.015), fill=LIGHT_GRAY)
    add_text(s, "推奨マーク", Inches(0.6), Inches(4.35), Inches(3), Inches(0.3),
             size=11, bold=True, color=NAVY)
    legends = [
        ("🔥", "最重要アップデート"),
        ("⚠️", "間違えやすいポイント"),
        ("💡", "副業で差がつくTips"),
        ("※要確認", "要裏取り"),
    ]
    lx = Inches(0.6)
    ly = Inches(4.75)
    for i, (mk, txt) in enumerate(legends):
        x = lx + (Inches(3.05)) * i
        add_text(s, mk, x, ly, Inches(0.7), Inches(0.4), size=14)
        add_text(s, txt, x + Inches(0.8), ly + Inches(0.05), Inches(2.3), Inches(0.3),
                 size=10, color=MID_GRAY)

    # "Why this deck" note
    add_rect(s, Inches(0.6), Inches(5.6), Inches(12.15), Inches(1.3), fill=ACCENT_BG, radius=True)
    add_text(s, "なぜ今このキャッチアップか",
             Inches(0.85), Inches(5.7), Inches(11), Inches(0.35),
             size=11, bold=True, color=ORANGE_DARK)
    add_text(s,
             "BSA戦略（Rapid HP制作+広告運用初月セット L3）で広告運用パートの商品力を強化する。\n"
             "同時に、副業として広告運用代行を再開できる実務ラインまで自分を戻す。",
             Inches(0.85), Inches(6.05), Inches(12), Inches(0.8),
             size=11, color=TEXT, line_spacing=1.4)

    footer(s, prs)


def slide_03_tldr(prs):
    s = blank(prs)
    page_frame(s, prs, "2.5年で何が変わったか — TL;DR", "4つの地殻変動", "P03 / 250")

    points = [
        ("01", "生成AIが\"ツール\"から\n\"運用レイヤー\"へ",
         "ChatGPT/Claude/Gemini/Sora…AIはツールから広告配信プラットフォームの内側に。クリエイティブ生成・入札最適化・アカウント構造提案まで AIが触る時代"),
        ("02", "Cookie廃止は撤回されたが\nシグナル減は戻らない",
         "2024/07/22 Google は Cookie 廃止計画を撤回。しかし ATT / 規制 / ブラウザ制限で計測シグナル減の趨勢は不変。CAPI / SGTM / Consent Mode v2 は依然必須"),
        ("03", "アカウント運用が\n\"ブラックボックス化\"",
         "PMax / Advantage+ / Smart+ に収束。広告セット分割・手動入札調整の時代は終わり、人間の仕事は「AIを導く / 検証設計 / 翻訳」へ"),
        ("04", "MMM が復活、\nIncrementality が標準装備",
         "ラストクリック計測崩壊の反動として Bayesian MMM（Meta Robyn / Google Meridian）が民主化。Meta は内製で Lift Test を標準提供"),
    ]
    top = Inches(1.5)
    h = Inches(2.6)
    w = Inches(5.85)
    gap = Inches(0.35)
    for i, (num, head, body) in enumerate(points):
        r = i // 2
        c = i % 2
        x = Inches(0.6) + (w + gap) * c
        y = top + (h + Inches(0.15)) * r
        add_rect(s, x, y, w, h, fill=WHITE, line=LIGHT_GRAY, radius=True)
        # big number
        add_text(s, num, x + Inches(0.3), y + Inches(0.2), Inches(1.2), Inches(0.9),
                 size=44, bold=True, color=ORANGE, font="Helvetica Neue")
        add_text(s, head, x + Inches(1.7), y + Inches(0.3), w - Inches(1.9), Inches(1.0),
                 size=16, bold=True, color=NAVY, line_spacing=1.2)
        add_text(s, body, x + Inches(0.3), y + Inches(1.45), w - Inches(0.6), h - Inches(1.55),
                 size=11, color=TEXT, line_spacing=1.4)

    footer(s, prs)


def slide_04_diffmap_a(prs):
    s = blank(prs)
    page_frame(s, prs, "差分マップ 2023-09 → 2026-04 (1/2)", "媒体・配信面", "P04 / 250")

    # Table header row
    headers = ["領域", "2023-09 当時", "2026-04 現在", "変化"]
    col_widths = [Inches(2.4), Inches(4.3), Inches(4.3), Inches(1.15)]
    col_x = [Inches(0.6)]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw)

    header_top = Inches(1.5)
    row_h = Inches(0.7)
    add_rect(s, Inches(0.6), header_top, Inches(12.15), row_h, fill=NAVY)
    for i, hd in enumerate(headers):
        align = PP_ALIGN.CENTER if i == 3 else PP_ALIGN.LEFT
        add_text(s, hd, col_x[i] + Inches(0.15), header_top + Inches(0.2),
                 col_widths[i] - Inches(0.15), Inches(0.4),
                 size=12, bold=True, color=WHITE, align=align)

    rows = [
        ("Google 検索広告",
         "BMM廃止後、Phrase/Exact で細かく制御。\nキーワード選定とマッチタイプ運用が肝",
         "Broad Match + Smart Bidding 推奨。\nAI Max for Search (2025/05) で自動拡張",
         "★★★"),
        ("Google 配信",
         "Smart Shopping / Local Campaigns 残存。\n手動キャンペーン分け運用",
         "Performance Max に収束。\nDemand Gen (Discovery統合)",
         "★★★"),
        ("Meta 配信",
         "ASC 初期、CBO 主流、\n広告セット分割が基本",
         "Advantage+ 全盛 (Shopping/Leads/App)。\nLattice / Andromeda で配信ML統合",
         "★★★"),
        ("TikTok",
         "日本B2C限定、手動運用。\nクリエイティブ職人的制作",
         "Smart+ (2024/10) / Symphony AI。\nShop Ads が本格拡大",
         "★★★"),
        ("Yahoo! / LINE",
         "Yahoo!広告単独、YDA 手動中心",
         "LINEヤフー統合後、LINE広告と\n連携前提の日本市場に",
         "★★"),
    ]
    y = header_top + row_h
    for i, row in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), y, Inches(12.15), Inches(0.95), fill=bg, line=LIGHT_GRAY)
        for j, cell in enumerate(row):
            align = PP_ALIGN.CENTER if j == 3 else PP_ALIGN.LEFT
            color = ORANGE if j == 3 else TEXT
            size = 14 if j == 3 else 10
            bold = j == 0 or j == 3
            add_text(s, cell, col_x[j] + Inches(0.15), y + Inches(0.12),
                     col_widths[j] - Inches(0.15), Inches(0.82),
                     size=size, bold=bold, color=NAVY if j == 0 else color,
                     align=align, line_spacing=1.3)
        y += Inches(0.95)

    # Legend bottom
    add_text(s, "★ = 変化の度合い (3段階)   |   出典: 各媒体公式発表 / Search Engine Land / 電通",
             Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
             size=9, color=SOFT_GRAY)

    footer(s, prs)


def slide_05_diffmap_b(prs):
    s = blank(prs)
    page_frame(s, prs, "差分マップ 2023-09 → 2026-04 (2/2)", "計測・クリエイティブ・運用職の役割", "P05 / 250")

    headers = ["領域", "2023-09 当時", "2026-04 現在", "変化"]
    col_widths = [Inches(2.4), Inches(4.3), Inches(4.3), Inches(1.15)]
    col_x = [Inches(0.6)]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw)

    header_top = Inches(1.5)
    row_h = Inches(0.7)
    add_rect(s, Inches(0.6), header_top, Inches(12.15), row_h, fill=NAVY)
    for i, hd in enumerate(headers):
        align = PP_ALIGN.CENTER if i == 3 else PP_ALIGN.LEFT
        add_text(s, hd, col_x[i] + Inches(0.15), header_top + Inches(0.2),
                 col_widths[i] - Inches(0.15), Inches(0.4),
                 size=12, bold=True, color=WHITE, align=align)

    rows = [
        ("計測基盤",
         "UA → GA4 移行中。\nタグ直接配置 / Cookieベース中心",
         "GA4 完全移行済。SGTM / CAPI 標準。\nConsent Mode v2 (2024/03~) 必須",
         "★★★"),
        ("クリエイティブ",
         "デザイナー制作。\n手動 A/B テストで勝ち本決定",
         "生成AI + Advantage+ Creative。\n動画は Sora / Veo / Runway",
         "★★★"),
        ("計測思想",
         "ラストクリック中心。\nCV重視のアトリビューション",
         "Incrementality / MMM / LTV へ回帰。\nMeta Robyn / Google Meridian 民主化",
         "★★"),
        ("運用職の役割",
         "配信設定 / 入札調整 /\nクリエイティブ指示",
         "AI を導く / 検証設計 / ビジネス翻訳。\n運用単純作業は AI へ",
         "★★★"),
        ("検索の位置づけ",
         "テキスト検索 = 購買導線の中心",
         "AI Overview (2024/05) / ChatGPT検索 /\nPerplexity で LLM経由流入が増加",
         "★★"),
    ]
    y = header_top + row_h
    for i, row in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), y, Inches(12.15), Inches(0.95), fill=bg, line=LIGHT_GRAY)
        for j, cell in enumerate(row):
            align = PP_ALIGN.CENTER if j == 3 else PP_ALIGN.LEFT
            color = ORANGE if j == 3 else TEXT
            size = 14 if j == 3 else 10
            bold = j == 0 or j == 3
            add_text(s, cell, col_x[j] + Inches(0.15), y + Inches(0.12),
                     col_widths[j] - Inches(0.15), Inches(0.82),
                     size=size, bold=bold, color=NAVY if j == 0 else color,
                     align=align, line_spacing=1.3)
        y += Inches(0.95)

    add_text(s, "→ 詳細は Part 3 (媒体別) / Part 4 (計測) / Part 5 (AI活用)",
             Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
             size=10, bold=True, color=ORANGE_DARK)

    footer(s, prs)


def slide_06_three_trends(prs):
    s = blank(prs)
    page_frame(s, prs, "業界3大潮流", "Three macro trends driving the landscape", "P06 / 250")

    trends = [
        ("01", "AI レイヤー化",
         "広告プラットフォーム各社が\n配信の中核に ML を据えた。",
         [
             "PMax / Advantage+ / Smart+ に収束",
             "人間は『AIを導く』役に",
             "代理店の省人化・内製化が加速",
         ]),
        ("02", "プライバシー持続",
         "Cookie 廃止は撤回されたが\nシグナル減少は不可逆。",
         [
             "CAPI / SGTM は標準装備",
             "Consent Mode v2 (2024/03~)",
             "Privacy Sandbox 終了 (2025/10) → CHIPS/FedCMのみ",
         ]),
        ("03", "購買導線の断片化",
         "検索→SNS→動画→AI検索へ\n導線が分散。",
         [
             "AI Overview (2024/05~)",
             "縦型動画の購買比率上昇",
             "GEO (Generative Engine Optimization) 胎動",
         ]),
    ]

    top = Inches(1.6)
    w = Inches(3.95)
    gap = Inches(0.15)
    h = Inches(5.1)
    for i, (num, title, lead, bullets) in enumerate(trends):
        x = Inches(0.6) + (w + gap) * i
        add_rect(s, x, top, w, h, fill=WHITE, line=LIGHT_GRAY, radius=True)
        # Top color band
        add_rect(s, x, top, w, Inches(0.15), fill=ORANGE)
        # Big number
        add_text(s, num, x + Inches(0.3), top + Inches(0.4), Inches(1.5), Inches(1.2),
                 size=52, bold=True, color=NAVY, font="Helvetica Neue")
        # Title
        add_text(s, title, x + Inches(0.3), top + Inches(1.55), w - Inches(0.6), Inches(0.6),
                 size=22, bold=True, color=NAVY)
        # Divider
        add_rect(s, x + Inches(0.3), top + Inches(2.15), w - Inches(0.6), Inches(0.015),
                 fill=ORANGE)
        # Lead
        add_text(s, lead, x + Inches(0.3), top + Inches(2.3), w - Inches(0.6), Inches(1.1),
                 size=12, color=TEXT, line_spacing=1.45)
        # Bullets
        by = top + Inches(3.45)
        for b in bullets:
            add_text(s, "▸", x + Inches(0.3), by, Inches(0.3), Inches(0.3),
                     size=11, color=ORANGE)
            add_text(s, b, x + Inches(0.55), by, w - Inches(0.9), Inches(0.5),
                     size=11, color=TEXT, line_spacing=1.3)
            by += Inches(0.45)

    footer(s, prs)


def slide_07_market_share(prs):
    s = blank(prs)
    page_frame(s, prs, "媒体シェア推移（日本市場）", "電通 日本の広告費 2024 (2025/02/27発表)", "P07 / 250")

    # Left: big number
    add_rect(s, Inches(0.6), Inches(1.5), Inches(4.5), Inches(5.3), fill=LIGHT, radius=True)
    add_text(s, "2024年 日本の総広告費", Inches(0.85), Inches(1.75), Inches(4.0), Inches(0.4),
             size=12, bold=True, color=NAVY)
    add_text(s, "7.67", Inches(0.85), Inches(2.25), Inches(3.5), Inches(1.2),
             size=72, bold=True, color=NAVY, font="Helvetica Neue")
    add_text(s, "兆円", Inches(3.2), Inches(2.85), Inches(1.0), Inches(0.6),
             size=22, bold=True, color=NAVY)
    add_text(s, "前年比 +4.9% ／ 3年連続で過去最高更新",
             Inches(0.85), Inches(3.65), Inches(4.0), Inches(0.35),
             size=11, color=MID_GRAY)
    # Sub stats
    add_rect(s, Inches(0.85), Inches(4.15), Inches(4.0), Inches(0.015), fill=LIGHT_GRAY)
    add_text(s, "うちインターネット広告費", Inches(0.85), Inches(4.3), Inches(4.0), Inches(0.35),
             size=11, bold=True, color=ORANGE_DARK)
    add_text(s, "3.65 兆円", Inches(0.85), Inches(4.65), Inches(4.0), Inches(0.6),
             size=30, bold=True, color=ORANGE, font="Helvetica Neue")
    add_text(s, "前年比 +9.6% ／ 全広告費の 47.6%",
             Inches(0.85), Inches(5.45), Inches(4.0), Inches(0.35),
             size=11, color=MID_GRAY)
    add_text(s, "うち広告媒体費 2.96 兆円 (+10.2%)",
             Inches(0.85), Inches(5.78), Inches(4.0), Inches(0.35),
             size=10, color=MID_GRAY)
    add_text(s, "出典 : 電通「2024年 日本の広告費」",
             Inches(0.85), Inches(6.4), Inches(4.0), Inches(0.35),
             size=9, color=SOFT_GRAY)

    # Right: bar comparison (totals by year)
    add_text(s, "内訳 : 媒体別広告費 (億円)",
             Inches(5.4), Inches(1.55), Inches(7.0), Inches(0.35),
             size=12, bold=True, color=NAVY)

    # Bars: ネット / テレビ / 新聞 / 雑誌 / ラジオ / OOH (billions yen approx from dentsu 2024)
    # Net: 36517, TV: ~17034, 新聞: 3351, 雑誌: 2072, ラジオ: 1063, OOH: 5996 (approx)
    bars = [
        ("インターネット", 36517, ORANGE),
        ("テレビメディア", 17034, NAVY),
        ("プロモーションメディア", 14959, NAVY),  # OOH+交通+DM+折込+フリーペーパー
        ("新聞", 3351, MID_GRAY),
        ("雑誌", 2072, MID_GRAY),
        ("ラジオ", 1063, MID_GRAY),
    ]
    max_val = max(v for _, v, _ in bars)
    bar_area_x = Inches(5.4)
    bar_area_y = Inches(2.05)
    label_w = Inches(2.3)
    bar_max_w = Inches(4.2)
    value_w = Inches(1.0)
    row_h = Inches(0.6)
    gap = Inches(0.1)
    for i, (label, val, color) in enumerate(bars):
        y = bar_area_y + (row_h + gap) * i
        add_text(s, label, bar_area_x, y + Inches(0.12), label_w, Inches(0.4),
                 size=11, bold=True, color=TEXT)
        bar_w = int(bar_max_w * (val / max_val))
        add_rect(s, bar_area_x + label_w, y + Inches(0.12), Emu(bar_w), Inches(0.35),
                 fill=color, radius=True)
        add_text(s, f"{val:,}", bar_area_x + label_w + Emu(bar_w) + Inches(0.1),
                 y + Inches(0.1), value_w, Inches(0.4),
                 size=11, color=MID_GRAY)

    # Callout
    add_rect(s, Inches(5.4), Inches(6.15), Inches(7.35), Inches(0.75), fill=ACCENT_BG, radius=True)
    add_text(s, "⚠", Inches(5.55), Inches(6.25), Inches(0.4), Inches(0.5),
             size=18, color=ORANGE_DARK)
    add_text(s,
             "運用型ネット広告がマス4媒体 (TV+新聞+雑誌+ラジオ) の合計を上回る構造が定着。\n"
             "リテールメディア / 縦型動画 / LINEヤフー合算シェアが今後の注目点。",
             Inches(6.0), Inches(6.2), Inches(6.7), Inches(0.7),
             size=10, color=TEXT, line_spacing=1.35)

    footer(s, prs)


def slide_08_keywords(prs):
    s = blank(prs)
    page_frame(s, prs, "2.5年で覚える10キーワード", "読む前の前提知識マップ", "P08 / 250")

    keywords = [
        ("Performance Max", "PMax",
         "Google全面統合型AIキャンペーン。検索以外はほぼこれに寄せる"),
        ("Advantage+", "ASC (Meta)",
         "Meta版の統合AI配信。Shopping / Leads / App。EC/D2Cの鉄板"),
        ("TikTok Smart+", "2024/10",
         "TikTok完全自動化キャンペーン。Web/Catalog/App/Lead"),
        ("AI Max for Search", "2025/05",
         "Google 検索広告のAI拡張層。DSAを置換予定"),
        ("CAPI", "Conversion API",
         "サーバーサイド計測の業界標準。Meta/TikTok/LINE/X対応"),
        ("Consent Mode v2", "2024/03~",
         "EU同意必須化。日本実務でも実装が定着"),
        ("SGTM", "Server-side GTM",
         "計測タグをサーバー経由で発火する標準構成"),
        ("MMM", "Media Mix Modeling",
         "Robyn / Meridian で民主化された統計モデル"),
        ("AI Overview / GEO", "2024/05~",
         "Google検索の AI 要約表示 → LLM経由流入の新潮流"),
        ("Agent型広告運用", "2025~",
         "Claude Computer Use / OpenAI Operator で管理画面を操作"),
    ]

    top = Inches(1.45)
    w = Inches(2.43)
    h = Inches(1.35)
    gap_x = Inches(0.06)
    gap_y = Inches(0.1)
    start_x = Inches(0.6)

    for i, (title, tag, desc) in enumerate(keywords):
        r = i // 5
        c = i % 5
        x = start_x + (w + gap_x) * c
        y = top + (h + gap_y) * r
        add_rect(s, x, y, w, h, fill=WHITE, line=LIGHT_GRAY, radius=True)
        # Number corner
        add_rect(s, x, y, Inches(0.08), h, fill=NAVY)
        # Number
        add_text(s, f"{i+1:02d}", x + Inches(0.2), y + Inches(0.12),
                 Inches(0.6), Inches(0.3),
                 size=11, bold=True, color=ORANGE, font="Helvetica Neue")
        # Tag
        add_text(s, tag, x + Inches(0.95), y + Inches(0.12), w - Inches(1.0), Inches(0.3),
                 size=9, color=SOFT_GRAY, align=PP_ALIGN.RIGHT)
        # Title
        add_text(s, title, x + Inches(0.2), y + Inches(0.45), w - Inches(0.3), Inches(0.4),
                 size=13, bold=True, color=NAVY)
        # Description
        add_text(s, desc, x + Inches(0.2), y + Inches(0.8), w - Inches(0.3), Inches(0.5),
                 size=9, color=TEXT, line_spacing=1.35)

    # Bottom note
    add_text(s,
             "各キーワードの深掘り : Part 3 (媒体固有) / Part 4 (計測) / Part 5 (AI活用) / 付録の用語集",
             Inches(0.6), Inches(6.7), Inches(12.15), Inches(0.3),
             size=10, color=MID_GRAY)

    footer(s, prs)


def slide_09_shortest_route(prs):
    s = blank(prs)
    page_frame(s, prs, "副業復帰の最短ルート", "Part 順ではなく効率順でなぞる", "P09 / 250")

    # 5 step chevron flow
    steps = [
        ("1", "Part 6", "鉄板運用メソッド\n2026年版", "15p / 30分",
         "現代の鉄板構成・入札戦略・\nLearning Phase 突破手順"),
        ("2", "Part 3-1,2", "Google + Meta\n重点読み", "47p / 60分",
         "PMax / Advantage+ の\n統制術を集中インストール"),
        ("3", "Part 4", "計測・データ\n基盤", "25p / 40分",
         "SGTM / CAPI /\nConsent Mode v2 押さえる"),
        ("4", "Part 5", "AI活用\nクリエイティブ", "40p / 60分",
         "生成AIツール3つ決める\n自分のワークフロー化"),
        ("5", "Part 7", "BSA L3\n実装ガイド", "10p / 20分",
         "30日プランで\n即クライアントに出せる状態"),
    ]

    top = Inches(1.6)
    box_w = Inches(2.3)
    arrow_w = Inches(0.15)
    gap = Inches(0.2)
    total = box_w * 5 + arrow_w * 4 + gap * 4
    start_x = Inches(0.6) + (Inches(12.15) - total) / 2

    for i, (n, part, title, time, desc) in enumerate(steps):
        x = start_x + (box_w + arrow_w + gap) * i
        # Box
        add_rect(s, x, top, box_w, Inches(4.0), fill=WHITE, line=LIGHT_GRAY, radius=True)
        # Step header
        add_rect(s, x, top, box_w, Inches(0.7), fill=NAVY, radius=True)
        add_text(s, f"STEP {n}", x + Inches(0.2), top + Inches(0.2),
                 box_w - Inches(0.4), Inches(0.35),
                 size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        # Part label
        add_text(s, part, x + Inches(0.2), top + Inches(0.85), box_w - Inches(0.4), Inches(0.35),
                 size=10, bold=True, color=ORANGE_DARK, align=PP_ALIGN.CENTER)
        # Title
        add_text(s, title, x + Inches(0.2), top + Inches(1.2), box_w - Inches(0.4), Inches(0.9),
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.3)
        # Time
        add_rect(s, x + Inches(0.4), top + Inches(2.1), box_w - Inches(0.8), Inches(0.015),
                 fill=LIGHT_GRAY)
        add_text(s, time, x + Inches(0.2), top + Inches(2.2), box_w - Inches(0.4), Inches(0.35),
                 size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # Desc
        add_text(s, desc, x + Inches(0.2), top + Inches(2.7), box_w - Inches(0.4), Inches(1.2),
                 size=10, color=TEXT, align=PP_ALIGN.CENTER, line_spacing=1.35)
        # Arrow
        if i < 4:
            arrow_x = x + box_w + gap / 2
            arrow_y = top + Inches(1.9)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       arrow_x, arrow_y, arrow_w + gap / 2, Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ORANGE
            arrow.line.fill.background()

    # Total
    add_rect(s, Inches(0.6), Inches(6.0), Inches(12.15), Inches(0.9), fill=NAVY, radius=True)
    add_text(s, "合計 約 3時間30分 で最低限の復帰ライン到達",
             Inches(0.6), Inches(6.17), Inches(12.15), Inches(0.55),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer(s, prs)


def slide_10_toc(prs):
    s = blank(prs)
    page_frame(s, prs, "目次", "Table of Contents", "P10 / 250")

    toc = [
        ("Part 1", "エグゼクティブサマリー", "10p", "P01"),
        ("Part 2", "市況・業界構造の変化", "20p", "P11"),
        ("Part 3", "媒体別キャッチアップ", "110p", "P31"),
        ("Part 4", "計測・データ基盤の激変", "25p", "P141"),
        ("Part 5", "AI活用事例集", "40p", "P166"),
        ("Part 6", "鉄板運用メソッド 2026年版", "15p", "P206"),
        ("Part 7", "BSA L3 実装ガイド", "10p", "P221"),
        ("Part 8", "付録 (用語・ツール・情報源)", "15p", "P231"),
    ]

    top = Inches(1.4)
    row_h = Inches(0.55)
    gap = Inches(0.07)

    for i, (part, title, pages, start) in enumerate(toc):
        y = top + (row_h + gap) * i
        # Part label pill
        add_rect(s, Inches(0.6), y, Inches(1.35), row_h, fill=NAVY, radius=True)
        add_text(s, part, Inches(0.6), y + Inches(0.15), Inches(1.35), Inches(0.35),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 font="Helvetica Neue")
        # Title
        add_rect(s, Inches(2.05), y, Inches(8.5), row_h, fill=LIGHT, radius=True)
        add_text(s, title, Inches(2.25), y + Inches(0.13), Inches(8.1), Inches(0.4),
                 size=14, bold=True, color=NAVY)
        # Pages
        add_rect(s, Inches(10.65), y, Inches(1.1), row_h, fill=WHITE, line=LIGHT_GRAY, radius=True)
        add_text(s, pages, Inches(10.65), y + Inches(0.15), Inches(1.1), Inches(0.35),
                 size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        # Start page
        add_text(s, f"→ {start}", Inches(11.85), y + Inches(0.15), Inches(0.9), Inches(0.35),
                 size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    # Bottom status
    add_rect(s, Inches(0.6), Inches(6.65), Inches(12.15), Inches(0.4), fill=ACCENT_BG, radius=True)
    add_text(s, "現在 Part 1 完了 | 総 250 ページ構成",
             Inches(0.6), Inches(6.72), Inches(12.15), Inches(0.3),
             size=11, bold=True, color=ORANGE_DARK, align=PP_ALIGN.CENTER)

    footer(s, prs)


# ---------- Main ----------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_cover(prs)
    slide_02_how_to_use(prs)
    slide_03_tldr(prs)
    slide_04_diffmap_a(prs)
    slide_05_diffmap_b(prs)
    slide_06_three_trends(prs)
    slide_07_market_share(prs)
    slide_08_keywords(prs)
    slide_09_shortest_route(prs)
    slide_10_toc(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
